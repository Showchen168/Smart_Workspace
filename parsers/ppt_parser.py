"""
PowerPoint 文件解析器
"""
from pptx import Presentation  # type: ignore[import-untyped]
from typing import Any, Dict, List, Optional


class PPTParser:
    """解析 PowerPoint (.pptx) 文件"""

    def __init__(self, file_path: str):
        self.file_path = file_path
        self.presentation: Any = None

    def parse(self) -> Dict[str, Any]:
        """
        解析 PPT 文件

        Returns:
            包含投影片內容的結構化字典
        """
        try:
            self.presentation = Presentation(self.file_path)

            content = {
                'file_name': self.file_path.split('\\')[-1],
                'file_type': 'pptx',
                'slides': self._extract_slides(),
                'full_text': self._extract_full_text()
            }

            return content

        except Exception as e:
            raise Exception(f"PPT 文件解析失敗: {str(e)}")

    def _extract_slides(self) -> List[Dict[str, Any]]:
        """提取每張投影片的內容"""
        slides = []
        assert self.presentation is not None

        for idx, slide in enumerate(self.presentation.slides, 1):
            slide_content = {
                'slide_number': idx,
                'title': self._extract_title(slide),
                'texts': self._extract_texts(slide),
                'notes': self._extract_notes(slide)
            }
            slides.append(slide_content)

        return slides

    def _extract_title(self, slide) -> str:
        """提取投影片標題"""
        if slide.shapes.title:
            return slide.shapes.title.text.strip()
        return ""

    def _extract_texts(self, slide) -> List[str]:
        """提取投影片中所有文字"""
        texts = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())

        return texts

    def _extract_notes(self, slide) -> str:
        """提取投影片備註"""
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame:
                return notes_frame.text.strip()
        return ""

    def _extract_full_text(self) -> str:
        """提取完整文字（用於 AI 分析）"""
        all_texts = []

        for slide_data in self._extract_slides():
            if slide_data['title']:
                all_texts.append(f"# {slide_data['title']}")
            all_texts.extend(slide_data['texts'])
            if slide_data['notes']:
                all_texts.append(f"備註: {slide_data['notes']}")

        return '\n\n'.join(all_texts)
